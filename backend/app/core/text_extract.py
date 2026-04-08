"""파일에서 텍스트를 추출하는 유틸리티.

지원 형식: PDF, PPTX, DOCX, TXT
"""
import io


def extract_text(content: bytes, mime_type: str) -> str:
    if mime_type == "application/pdf":
        return _from_pdf(content)
    if mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        return _from_pptx(content)
    if mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        return _from_docx(content)
    if mime_type == "text/plain":
        return content.decode("utf-8", errors="ignore")
    raise ValueError(f"지원하지 않는 MIME 타입: {mime_type}")


def _from_pdf(content: bytes) -> str:
    import pdfplumber

    text_parts = []
    with pdfplumber.open(io.BytesIO(content)) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                text_parts.append(text)
    return "\n\n".join(text_parts)


def _from_pptx(content: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(content))
    text_parts = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = " ".join(run.text for run in para.runs).strip()
                    if line:
                        slide_texts.append(line)
        if slide_texts:
            text_parts.append(f"[슬라이드 {slide_num}]\n" + "\n".join(slide_texts))
    return "\n\n".join(text_parts)


def _from_docx(content: bytes) -> str:
    from docx import Document

    doc = Document(io.BytesIO(content))
    return "\n".join(para.text for para in doc.paragraphs if para.text.strip())
